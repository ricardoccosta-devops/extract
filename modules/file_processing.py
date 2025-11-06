"""Módulo de processamento e extração de texto de arquivos."""

import io
from dataclasses import dataclass
from datetime import datetime
from typing import Optional

import structlog
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

logger = structlog.get_logger()


@dataclass
class ProcessedDocument:
    """Estrutura de dados para documentos processados."""

    content: str
    metadata: dict
    file_name: str
    file_size: int
    file_type: str
    processed_at: datetime


class FileProcessingService:
    """Serviço para processamento e extração de texto de arquivos."""

    SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".docx", ".doc", ".pptx", ".ppt"}

    def __init__(self):
        """Inicializa o serviço de processamento."""
        self.logger = logger

    def process_file(
        self, file_content: bytes, file_name: str
    ) -> ProcessedDocument:
        """
        Processa um arquivo e extrai seu conteúdo e metadados.

        Args:
            file_content: Conteúdo do arquivo em bytes
            file_name: Nome do arquivo

        Returns:
            ProcessedDocument com conteúdo e metadados

        Raises:
            ValueError: Se o formato do arquivo não for suportado
        """
        self.logger.info("Iniciando processamento de arquivo", file_name=file_name)

        file_extension = self._get_file_extension(file_name)
        if file_extension not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Formato não suportado: {file_extension}. "
                f"Formatos suportados: {', '.join(self.SUPPORTED_EXTENSIONS)}"
            )

        try:
            content = self._extract_text(file_content, file_extension)
            metadata = self._extract_metadata(file_content, file_extension, file_name)

            processed_doc = ProcessedDocument(
                content=content,
                metadata=metadata,
                file_name=file_name,
                file_size=len(file_content),
                file_type=file_extension,
                processed_at=datetime.now(),
            )

            self.logger.info(
                "Arquivo processado com sucesso",
                file_name=file_name,
                content_length=len(content),
            )
            return processed_doc

        except Exception as e:
            self.logger.error("Erro ao processar arquivo", error=str(e), file_name=file_name)
            raise

    def _get_file_extension(self, file_name: str) -> str:
        """Extrai a extensão do arquivo."""
        return "." + file_name.split(".")[-1].lower()

    def _extract_text(self, file_content: bytes, file_extension: str) -> str:
        """Extrai texto do arquivo baseado na extensão."""
        if file_extension == ".pdf":
            return self._extract_from_pdf(file_content)
        elif file_extension == ".txt":
            return self._extract_from_txt(file_content)
        elif file_extension in [".docx", ".doc"]:
            return self._extract_from_docx(file_content)
        elif file_extension in [".pptx", ".ppt"]:
            return self._extract_from_pptx(file_content)
        else:
            raise ValueError(f"Extrator não implementado para: {file_extension}")

    def _extract_from_pdf(self, file_content: bytes) -> str:
        """Extrai texto de um arquivo PDF."""
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PdfReader(pdf_file)
            text_parts = []

            for page_num, page in enumerate(pdf_reader.pages, 1):
                try:
                    text = page.extract_text()
                    if text.strip():
                        text_parts.append(f"--- Página {page_num} ---\n{text}")
                except Exception as e:
                    self.logger.warning(
                        "Erro ao extrair texto da página",
                        page=page_num,
                        error=str(e),
                    )

            return "\n\n".join(text_parts)

        except Exception as e:
            self.logger.error("Erro ao processar PDF", error=str(e))
            raise ValueError(f"Erro ao processar PDF: {e}")

    def _extract_from_txt(self, file_content: bytes) -> str:
        """Extrai texto de um arquivo TXT."""
        try:
            # Tenta diferentes encodings
            encodings = ["utf-8", "latin-1", "cp1252"]
            for encoding in encodings:
                try:
                    return file_content.decode(encoding)
                except UnicodeDecodeError:
                    continue
            raise ValueError("Não foi possível decodificar o arquivo TXT")
        except Exception as e:
            self.logger.error("Erro ao processar TXT", error=str(e))
            raise ValueError(f"Erro ao processar TXT: {e}")

    def _extract_from_docx(self, file_content: bytes) -> str:
        """Extrai texto de um arquivo DOCX."""
        try:
            doc_file = io.BytesIO(file_content)
            doc = Document(doc_file)
            text_parts = []

            # Extrai parágrafos
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)

            # Extrai tabelas
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                if table_text:
                    text_parts.append("\n".join(table_text))

            return "\n\n".join(text_parts)

        except Exception as e:
            self.logger.error("Erro ao processar DOCX", error=str(e))
            raise ValueError(f"Erro ao processar DOCX: {e}")

    def _extract_from_pptx(self, file_content: bytes) -> str:
        """Extrai texto de um arquivo PPTX."""
        try:
            pptx_file = io.BytesIO(file_content)
            prs = Presentation(pptx_file)
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if len(slide_text) > 1:
                    text_parts.append("\n".join(slide_text))

            return "\n\n".join(text_parts)

        except Exception as e:
            self.logger.error("Erro ao processar PPTX", error=str(e))
            raise ValueError(f"Erro ao processar PPTX: {e}")

    def _extract_metadata(
        self, file_content: bytes, file_extension: str, file_name: str
    ) -> dict:
        """Extrai metadados do arquivo."""
        metadata = {
            "file_name": file_name,
            "file_extension": file_extension,
            "file_size_bytes": len(file_content),
            "file_size_mb": round(len(file_content) / (1024 * 1024), 2),
        }

        try:
            if file_extension == ".pdf":
                pdf_file = io.BytesIO(file_content)
                pdf_reader = PdfReader(pdf_file)
                metadata.update(
                    {
                        "num_pages": len(pdf_reader.pages),
                        "is_encrypted": pdf_reader.is_encrypted,
                    }
                )
                if pdf_reader.metadata:
                    metadata["pdf_metadata"] = {
                        k: str(v) for k, v in pdf_reader.metadata.items()
                    }

            elif file_extension in [".docx", ".doc"]:
                doc_file = io.BytesIO(file_content)
                doc = Document(doc_file)
                metadata.update(
                    {
                        "num_paragraphs": len(doc.paragraphs),
                        "num_tables": len(doc.tables),
                    }
                )
                if doc.core_properties:
                    metadata["docx_metadata"] = {
                        "title": doc.core_properties.title,
                        "author": doc.core_properties.author,
                        "created": str(doc.core_properties.created),
                        "modified": str(doc.core_properties.modified),
                    }

            elif file_extension in [".pptx", ".ppt"]:
                pptx_file = io.BytesIO(file_content)
                prs = Presentation(pptx_file)
                metadata.update({"num_slides": len(prs.slides)})

        except Exception as e:
            self.logger.warning("Erro ao extrair metadados", error=str(e))

        return metadata

